"""
PPTX Validation Script
Checks for common layout and content issues in PowerPoint presentations.

Checks performed:
1. Text boxes extending beyond slide boundaries (10" x 5.625")
2. Overlapping text boxes
3. Text that might be cut off (too much text for allocated space)
4. Missing content / empty placeholders
5. Elements positioned too close to edges
6. Font size consistency
7. Slide structure summary
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

# Standard slide dimensions (widescreen 16:9)
SLIDE_WIDTH_INCHES = 13.333  # Standard widescreen
SLIDE_HEIGHT_INCHES = 7.5
EDGE_MARGIN_INCHES = 0.25  # Minimum safe margin from edges

# Convert to EMU for precision
SLIDE_WIDTH_EMU = Inches(SLIDE_WIDTH_INCHES)
SLIDE_HEIGHT_EMU = Inches(SLIDE_HEIGHT_INCHES)
EDGE_MARGIN_EMU = Inches(EDGE_MARGIN_INCHES)


def emu_to_inches(emu):
    """Convert EMU to inches."""
    if emu is None:
        return None
    return round(emu / 914400, 3)


def estimate_text_lines(text, font_size_pt, box_width_inches):
    """Rough estimate of how many lines of text will be needed."""
    if not text or not font_size_pt or not box_width_inches:
        return 0
    # Approximate: average char width ~ 0.6 * font_size in points, 1 point = 1/72 inch
    avg_char_width_inches = 0.5 * font_size_pt / 72
    if avg_char_width_inches == 0:
        return 0
    chars_per_line = max(1, int(box_width_inches / avg_char_width_inches))
    lines = 0
    for paragraph in text.split('\n'):
        if not paragraph:
            lines += 1
        else:
            lines += max(1, -(-len(paragraph) // chars_per_line))  # Ceiling division
    return lines


def get_shape_rect(shape):
    """Get shape bounding rectangle in EMU."""
    return {
        'left': shape.left,
        'top': shape.top,
        'right': shape.left + shape.width,
        'bottom': shape.top + shape.height,
        'width': shape.width,
        'height': shape.height,
    }


def rects_overlap(r1, r2):
    """Check if two rectangles overlap."""
    if r1['right'] <= r2['left'] or r2['right'] <= r1['left']:
        return False
    if r1['bottom'] <= r2['top'] or r2['bottom'] <= r1['top']:
        return False
    return True


def overlap_area(r1, r2):
    """Calculate overlap area in square inches."""
    x_overlap = max(0, min(r1['right'], r2['right']) - max(r1['left'], r2['left']))
    y_overlap = max(0, min(r1['bottom'], r2['bottom']) - max(r1['top'], r2['top']))
    area_emu2 = x_overlap * y_overlap
    return round(area_emu2 / (914400 * 914400), 4)  # Convert to square inches


def validate_presentation(filepath):
    """Validate a PPTX file and return issues found."""
    issues = []
    warnings = []
    info = []

    if not os.path.exists(filepath):
        return {"error": f"File not found: {filepath}"}

    prs = Presentation(filepath)

    # Get actual slide dimensions
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_w_in = emu_to_inches(slide_w)
    slide_h_in = emu_to_inches(slide_h)

    info.append(f"Slide dimensions: {slide_w_in}\" x {slide_h_in}\" ({slide_w} x {slide_h} EMU)")
    info.append(f"Number of slides: {len(prs.slides)}")

    # Use actual slide dimensions for checks
    actual_slide_w = slide_w
    actual_slide_h = slide_h

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_issues = []
        slide_warnings = []
        text_shapes = []

        # Collect all shapes with text
        all_shapes = list(slide.shapes)

        info.append(f"\n{'='*60}")
        info.append(f"SLIDE {slide_idx}: {len(all_shapes)} shapes")

        # Check slide layout name
        if slide.slide_layout:
            info.append(f"  Layout: {slide.slide_layout.name}")

        for shape in all_shapes:
            shape_name = shape.name or "(unnamed)"
            rect = get_shape_rect(shape)

            shape_info = {
                'name': shape_name,
                'type': str(shape.shape_type),
                'left': emu_to_inches(shape.left),
                'top': emu_to_inches(shape.top),
                'width': emu_to_inches(shape.width),
                'height': emu_to_inches(shape.height),
                'right': emu_to_inches(rect['right']),
                'bottom': emu_to_inches(rect['bottom']),
            }

            info.append(f"  Shape: '{shape_name}' | Type: {shape.shape_type} | "
                       f"Pos: ({shape_info['left']}\", {shape_info['top']}\") | "
                       f"Size: {shape_info['width']}\" x {shape_info['height']}\"")

            # CHECK 1: Extends beyond slide boundaries
            if rect['right'] > actual_slide_w + Inches(0.01):  # Small tolerance
                overshoot = emu_to_inches(rect['right'] - actual_slide_w)
                slide_issues.append(
                    f"BOUNDARY: '{shape_name}' extends {overshoot}\" beyond RIGHT edge "
                    f"(right={shape_info['right']}\", slide width={emu_to_inches(actual_slide_w)}\")"
                )

            if rect['bottom'] > actual_slide_h + Inches(0.01):
                overshoot = emu_to_inches(rect['bottom'] - actual_slide_h)
                slide_issues.append(
                    f"BOUNDARY: '{shape_name}' extends {overshoot}\" beyond BOTTOM edge "
                    f"(bottom={shape_info['bottom']}\", slide height={emu_to_inches(actual_slide_h)}\")"
                )

            if rect['left'] < -Inches(0.01):
                slide_issues.append(
                    f"BOUNDARY: '{shape_name}' extends beyond LEFT edge "
                    f"(left={shape_info['left']}\")"
                )

            if rect['top'] < -Inches(0.01):
                slide_issues.append(
                    f"BOUNDARY: '{shape_name}' extends beyond TOP edge "
                    f"(top={shape_info['top']}\")"
                )

            # CHECK 5: Too close to edges
            if rect['left'] >= 0 and rect['left'] < EDGE_MARGIN_EMU and shape_info['width'] > 0.5:
                slide_warnings.append(
                    f"EDGE: '{shape_name}' is only {shape_info['left']}\" from LEFT edge"
                )

            if rect['top'] >= 0 and rect['top'] < EDGE_MARGIN_EMU and shape_info['height'] > 0.5:
                slide_warnings.append(
                    f"EDGE: '{shape_name}' is only {shape_info['top']}\" from TOP edge"
                )

            right_margin = emu_to_inches(actual_slide_w - rect['right'])
            if right_margin >= 0 and right_margin < EDGE_MARGIN_INCHES and shape_info['width'] > 0.5:
                slide_warnings.append(
                    f"EDGE: '{shape_name}' is only {right_margin}\" from RIGHT edge"
                )

            bottom_margin = emu_to_inches(actual_slide_h - rect['bottom'])
            if bottom_margin >= 0 and bottom_margin < EDGE_MARGIN_INCHES and shape_info['height'] > 0.5:
                slide_warnings.append(
                    f"EDGE: '{shape_name}' is only {bottom_margin}\" from BOTTOM edge"
                )

            # CHECK 3 & 4: Text content analysis
            if shape.has_text_frame:
                tf = shape.text_frame
                full_text = tf.text.strip()

                if not full_text:
                    slide_warnings.append(
                        f"EMPTY: '{shape_name}' has a text frame but no text content"
                    )
                else:
                    # Get text details
                    text_preview = full_text[:80] + "..." if len(full_text) > 80 else full_text
                    info.append(f"    Text: \"{text_preview}\"")

                    # Collect font sizes
                    font_sizes = set()
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_sizes.add(run.font.size)

                    if font_sizes:
                        sizes_pt = [round(s / 12700, 1) for s in font_sizes]
                        info.append(f"    Font sizes: {sorted(sizes_pt)} pt")

                        # Estimate if text fits
                        min_font = min(font_sizes) / 12700  # to points
                        box_w_in = emu_to_inches(shape.width)
                        box_h_in = emu_to_inches(shape.height)

                        est_lines = estimate_text_lines(full_text, min_font, box_w_in)
                        line_height_in = min_font * 1.2 / 72  # Approximate line height
                        max_lines = int(box_h_in / line_height_in) if line_height_in > 0 else 999

                        if est_lines > max_lines * 1.1:  # 10% tolerance
                            slide_warnings.append(
                                f"OVERFLOW: '{shape_name}' may have text overflow "
                                f"(~{est_lines} lines estimated, ~{max_lines} lines fit, "
                                f"font={min_font}pt, box={box_w_in}\"x{box_h_in}\")"
                            )

                    # Check for very small font sizes
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size < Pt(8):
                                size_pt = round(run.font.size / 12700, 1)
                                slide_warnings.append(
                                    f"SMALL_FONT: '{shape_name}' has text at {size_pt}pt "
                                    f"(may be hard to read)"
                                )
                                break

                text_shapes.append((shape, rect))

        # CHECK 2: Overlapping text boxes
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                s1, r1 = text_shapes[i]
                s2, r2 = text_shapes[j]

                if rects_overlap(r1, r2):
                    area = overlap_area(r1, r2)
                    if area > 0.1:  # More than 0.1 sq inch overlap
                        slide_warnings.append(
                            f"OVERLAP: '{s1.name}' and '{s2.name}' overlap "
                            f"by ~{area} sq inches"
                        )

        # Report slide results
        if slide_issues:
            for issue in slide_issues:
                issues.append(f"Slide {slide_idx} - {issue}")
        if slide_warnings:
            for warning in slide_warnings:
                warnings.append(f"Slide {slide_idx} - {warning}")

    return {
        'info': info,
        'issues': issues,
        'warnings': warnings,
    }


def print_report(filepath, results):
    """Print a formatted report."""
    basename = os.path.basename(filepath)

    print(f"\n{'#'*70}")
    print(f"# VALIDATION REPORT: {basename}")
    print(f"{'#'*70}")

    if 'error' in results:
        print(f"\nERROR: {results['error']}")
        return

    # Info section
    print("\n--- SLIDE STRUCTURE ---")
    for line in results['info']:
        print(line)

    # Issues (critical)
    print(f"\n--- CRITICAL ISSUES ({len(results['issues'])}) ---")
    if results['issues']:
        for issue in results['issues']:
            print(f"  [ERROR] {issue}")
    else:
        print("  None found.")

    # Warnings
    print(f"\n--- WARNINGS ({len(results['warnings'])}) ---")
    if results['warnings']:
        for warning in results['warnings']:
            print(f"  [WARN] {warning}")
    else:
        print("  None found.")

    # Summary
    total_problems = len(results['issues']) + len(results['warnings'])
    print(f"\n--- SUMMARY ---")
    print(f"  Critical issues: {len(results['issues'])}")
    print(f"  Warnings: {len(results['warnings'])}")
    if total_problems == 0:
        print("  STATUS: PASS - No issues detected")
    elif results['issues']:
        print("  STATUS: FAIL - Critical issues found that need fixing")
    else:
        print("  STATUS: REVIEW - Warnings found, manual review recommended")


if __name__ == '__main__':
    files = [
        r"C:\Users\Philippe\Documents\CLAUDECODE\ELECTIONS-LAFLECHE\docs\synthese-decideurs.pptx",
        r"C:\Users\Philippe\Documents\CLAUDECODE\ELECTIONS-LAFLECHE\docs\synthese-militants.pptx",
    ]

    for filepath in files:
        results = validate_presentation(filepath)
        print_report(filepath, results)

    print("\n" + "="*70)
    print("Validation complete.")
